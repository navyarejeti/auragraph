"""
backend/agents/slide_images.py
───────────────────────────────
Extract embedded images from PDF slides and PPTX files.

Fixes in v5:
  D1 – PPTX pixel calc: `shape.width or 1` → 0px. Now uses proper cm→px
       conversion with a minimum 10px fallback before the size filter.
  D2 – `len(doc)` called AFTER `doc.close()`. Fixed by saving page_count first.
  D3 – MAX_IMAGES_PER_FILE with no per-page budget. Now 4 images per page max,
       still honouring the total cap so later slides always get a fair share.
  D4 – save_images() never cleared stale images on re-upload. Now accepts a
       `clear_existing` flag; upload handler passes True for slide images.
"""

from __future__ import annotations

import hashlib
import io
import logging
import os
import shutil
from pathlib import Path
from dataclasses import dataclass, field

logger = logging.getLogger(__name__)

MAX_IMAGES_PER_FILE  = 20   # total cap per file
MAX_IMAGES_PER_PAGE  =  4   # FIX D3: per-page budget
MIN_DIM_PX           = 80
MIN_AREA_PX          = 10_000
MIN_BYTES            = 4_000


@dataclass
class ExtractedImage:
    img_id:       str
    data:         bytes
    mime:         str
    source_label: str
    width:        int
    height:       int
    description:  str = ""


# ── PDF extraction via PyMuPDF ─────────────────────────────────────────────

def extract_images_from_pdf(file_bytes: bytes) -> list[ExtractedImage]:
    """Extract significant embedded images from a PDF."""
    results: list[ExtractedImage] = []
    seen_hashes: set[str] = set()
    counter = 1

    try:
        import fitz
    except ImportError:
        logger.warning("slide_images: PyMuPDF not installed — PDF image extraction unavailable")
        return []

    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
    except Exception as e:
        logger.warning("slide_images: fitz failed to open PDF: %s", e)
        return []

    # FIX D2: record page count BEFORE closing doc
    page_count = len(doc)

    for page_num, page in enumerate(doc, start=1):
        if len(results) >= MAX_IMAGES_PER_FILE:
            break

        page_count_this = 0  # FIX D3: per-page counter
        try:
            img_list = page.get_images(full=True)
        except Exception:
            continue

        for img_info in img_list:
            if len(results) >= MAX_IMAGES_PER_FILE:
                break
            # FIX D3: respect per-page budget
            if page_count_this >= MAX_IMAGES_PER_PAGE:
                break
            try:
                xref = img_info[0]
                w    = img_info[2]
                h    = img_info[3]

                if w < MIN_DIM_PX or h < MIN_DIM_PX:
                    continue
                if w * h < MIN_AREA_PX:
                    continue

                base_img   = doc.extract_image(xref)
                data       = base_img["image"]
                mime_short = base_img.get("ext", "png").lower()
                mime       = f"image/{mime_short}".replace("image/jpg", "image/jpeg")

                if len(data) < MIN_BYTES:
                    continue

                h_str = hashlib.md5(data).hexdigest()
                if h_str in seen_hashes:
                    continue
                seen_hashes.add(h_str)

                img_id = f"img_{counter:03d}"
                counter += 1
                page_count_this += 1
                results.append(ExtractedImage(
                    img_id=img_id, data=data, mime=mime,
                    source_label=f"Page {page_num}", width=w, height=h,
                ))
            except Exception as e:
                logger.debug("slide_images: skipping PDF img xref %s: %s", img_info[0], e)

    # FIX D2: close first, then use saved page_count
    doc.close()
    logger.info("slide_images: extracted %d images from PDF (%d pages)", len(results), page_count)
    return results


# ── PPTX extraction via python-pptx ──────────────────────────────────────────

def extract_images_from_pptx(file_bytes: bytes) -> list[ExtractedImage]:
    """Extract significant embedded images from a PPTX file."""
    results: list[ExtractedImage] = []
    seen_hashes: set[str] = set()
    counter = 1

    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("slide_images: python-pptx not installed")
        return []

    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception as e:
        logger.warning("slide_images: python-pptx failed to open file: %s", e)
        return []

    for slide_num, slide in enumerate(prs.slides, start=1):
        if len(results) >= MAX_IMAGES_PER_FILE:
            break

        page_count_this = 0  # FIX D3: per-slide budget
        for shape in slide.shapes:
            if len(results) >= MAX_IMAGES_PER_FILE:
                break
            if page_count_this >= MAX_IMAGES_PER_PAGE:  # FIX D3
                break
            try:
                if getattr(shape, 'shape_type', None) != 13:  # MSO_SHAPE_TYPE.PICTURE
                    continue
                img  = shape.image
                data = img.blob
                if len(data) < MIN_BYTES:
                    continue

                # FIX D1: correct EMU→px conversion, guard against zero/tiny EMU
                # 914400 EMU = 1 inch; assume 96 DPI screen rendering
                w_emu = shape.width  or 0
                h_emu = shape.height or 0
                # Minimum 10px so we can still apply the MIN_DIM_PX filter correctly
                w_px = max(10, int(w_emu / 914400 * 96)) if w_emu > 0 else 10
                h_px = max(10, int(h_emu / 914400 * 96)) if h_emu > 0 else 10

                if w_px < MIN_DIM_PX or h_px < MIN_DIM_PX:
                    continue

                ext  = img.ext.lower()
                mime = f"image/{ext}".replace("image/jpg", "image/jpeg")

                h_str = hashlib.md5(data).hexdigest()
                if h_str in seen_hashes:
                    continue
                seen_hashes.add(h_str)

                img_id = f"img_{counter:03d}"
                counter += 1
                page_count_this += 1
                results.append(ExtractedImage(
                    img_id=img_id, data=data, mime=mime,
                    source_label=f"Slide {slide_num}", width=w_px, height=h_px,
                ))
            except Exception as e:
                logger.debug("slide_images: skipping PPTX shape: %s", e)

    logger.info("slide_images: extracted %d images from PPTX (%d slides)",
                len(results), len(prs.slides))
    return results


# ── Unified entry point ────────────────────────────────────────────────────────

def extract_images_from_file(file_bytes: bytes, filename: str) -> list[ExtractedImage]:
    from pathlib import Path
    ext = Path(filename.lower()).suffix
    if ext in {'.pptx', '.ppt'}:
        return extract_images_from_pptx(file_bytes)
    if ext == '.pdf':
        return extract_images_from_pdf(file_bytes)
    return []


# ── Image store (persistent filesystem) ───────────────────────────────────────

_DEFAULT_IMG_ROOT = str(Path(__file__).resolve().parents[1] / "image_store")
_IMG_ROOT = os.environ.get("AURAGRAPH_IMAGE_ROOT", _DEFAULT_IMG_ROOT)


def save_images(notebook_id: str, images: list[ExtractedImage],
                clear_existing: bool = False) -> None:
    """
    Save image bytes to disk so the API endpoint can serve them.

    FIX D4: `clear_existing=True` wipes the folder first so stale images
    from a previous upload of the same notebook do not linger on disk.
    Pass True for slide images (re-upload replaces them).
    Pass False for textbook images being appended in the same request.
    """
    folder = os.path.join(_IMG_ROOT, notebook_id)
    # FIX D4: clear folder before writing slide images
    if clear_existing and os.path.isdir(folder):
        shutil.rmtree(folder)
    os.makedirs(folder, exist_ok=True)
    for img in images:
        ext  = img.mime.split("/")[-1].replace("jpeg", "jpg")
        path = os.path.join(folder, f"{img.img_id}.{ext}")
        with open(path, "wb") as f:
            f.write(img.data)
    logger.info("slide_images: saved %d images for notebook %s (clear=%s)",
                len(images), notebook_id, clear_existing)


def get_image_path(notebook_id: str, img_id_with_ext: str) -> str | None:
    folder = os.path.join(_IMG_ROOT, notebook_id)
    path   = os.path.join(folder, img_id_with_ext)
    return path if os.path.exists(path) else None
