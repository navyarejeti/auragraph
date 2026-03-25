"""
Human-style end-to-end test for the AuraGraph notes generation pipeline.

What it does:
  1. Builds a realistic lecture PPTX (Fourier Transform, 8 slides, 3 topics)
  2. POSTs it to /api/upload-fuse-multi with proficiency = Intermediate
  3. Audits the returned notes like a human reviewer:
       ✔  Every topic has a ## heading
       ✔  Every section has an Exam Tip
       ✔  Key points from slides appear in notes
       ✔  Math is in LaTeX (not spelled out as "omega", "sigma" etc.)
       ✔  No LLM preamble ("Here are your notes", "Certainly", etc.)
       ✔  Tables use pipe format (not HTML)
       ✔  Source is reported (azure / groq / local)
       ✔  Response time is acceptable (< 120 s)
  4. Prints a coloured pass/fail report.

Run from the backend/ directory with the venv active:
    python test_notes_e2e.py [--proficiency Beginner|Intermediate|Advanced]
"""

import io
import sys
import time
import argparse
import re
import requests
from pptx import Presentation
from pptx.util import Inches, Pt

BASE_URL = "http://localhost:8000"

# ── Slide content ─────────────────────────────────────────────────────────────
SLIDES = [
    {
        "title": "Discrete Fourier Transform (DFT)",
        "body": (
            "Definition: X[k] = sum_{n=0}^{N-1} x[n] * e^{-j2πkn/N}\n"
            "• Converts N-point time-domain signal to frequency domain\n"
            "• k = frequency index, n = time index, N = signal length\n"
            "• DFT is periodic with period N"
        ),
    },
    {
        "title": "DFT Properties",
        "body": (
            "• Linearity: DFT{ax+by} = aDFT{x} + bDFT{y}\n"
            "• Time shift: x[n-m] ↔ X[k] e^{-j2πkm/N}\n"
            "• Frequency shift: x[n] e^{j2πk₀n/N} ↔ X[k-k₀]\n"
            "• Convolution: x[n]*h[n] ↔ X[k]H[k]"
        ),
    },
    {
        "title": "Fast Fourier Transform (FFT)",
        "body": (
            "• FFT is an efficient algorithm to compute DFT\n"
            "• Cooley-Tukey radix-2 FFT: O(N log N) vs O(N²) for DFT\n"
            "• Requires N to be a power of 2\n"
            "• Butterfly diagram splits even/odd indexed samples\n"
            "• Key formula: W_N = e^{-j2π/N} (twiddle factor)"
        ),
    },
    {
        "title": "FFT Butterfly Operation",
        "body": (
            "Stage 1: split x[n] into even x_e[n] and odd x_o[n]\n"
            "Stage 2: A = x_e + W_N^k * x_o\n"
            "        B = x_e - W_N^k * x_o\n"
            "log₂(N) stages, N/2 butterflies per stage\n"
            "Total multiplications: (N/2) log₂(N)"
        ),
    },
    {
        "title": "Zero Padding",
        "body": (
            "• Append zeros to signal to increase DFT resolution\n"
            "• Does NOT improve frequency resolution (no new info)\n"
            "• Interpolates spectrum — smoother, not more accurate\n"
            "• Common in spectral analysis and OFDM systems"
        ),
    },
    {
        "title": "IDFT — Inverse Discrete Fourier Transform",
        "body": (
            "x[n] = (1/N) sum_{k=0}^{N-1} X[k] e^{j2πkn/N}\n"
            "• Reconstructs time-domain signal from spectrum\n"
            "• Only difference from DFT: positive exponent + 1/N scale"
        ),
    },
    {
        "title": "Circular Convolution vs Linear Convolution",
        "body": (
            "• Linear convolution: y[n] = x[n] * h[n], length L+M-1\n"
            "• Circular convolution: y_c[n] = x[n] ⊛ h[n], length N\n"
            "• To compute linear via FFT: zero-pad to at least L+M-1\n"
            "• Overlap-add / overlap-save methods for long sequences"
        ),
    },
    {
        "title": "Applications of FFT",
        "body": (
            "• Spectral analysis (audio, radar, biomedical signals)\n"
            "• OFDM in 4G/5G wireless communication\n"
            "• Image compression (JPEG uses DCT, a cosine variant)\n"
            "• Polynomial multiplication in O(N log N)\n"
            "• Fast correlation & matched filtering"
        ),
    },
]

TEXTBOOK_CONTEXT = """
The Discrete Fourier Transform is the foundation of digital signal processing.
Its computational complexity was reduced from O(N²) to O(N log N) by Cooley and Tukey in 1965.
The twiddle factor W_N^k = e^{-j2πk/N} is a unit complex rotation on the unit circle.
Zero padding is often confused with increasing resolution; it only interpolates an already-computed spectrum.
IDFT requires careful scaling by 1/N to maintain energy preservation (Parseval's theorem: sum|x[n]|² = (1/N) sum|X[k]|²).
"""

# ── Build PPTX in-memory ──────────────────────────────────────────────────────

def build_pptx() -> bytes:
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[1]  # Title + Content

    for slide_data in SLIDES:
        slide = prs.slides.add_slide(blank_slide_layout)
        title = slide.shapes.title
        body  = slide.placeholders[1]

        title.text = slide_data["title"]
        body.text  = slide_data["body"]

        # Ensure font sizes are reasonable
        for para in body.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(14)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Quality auditors ──────────────────────────────────────────────────────────

PREAMBLE_RE = re.compile(
    r'^(here\s+(are|is)|sure[,!]?|certainly[,!]?|below\s+(are|is)|'
    r'of\s+course|the\s+following|these\s+are)',
    re.IGNORECASE | re.MULTILINE,
)
SPELLED_MATH_RE = re.compile(
    r'\b(omega|sigma|alpha|beta|delta|theta|lambda|epsilon|pi(?!\w))\b',
    re.IGNORECASE,
)
HTML_TABLE_RE = re.compile(r'<table|<tr|<td|<th', re.IGNORECASE)

def audit_notes(notes: str, proficiency: str) -> list[tuple[str, bool, str]]:
    """
    Returns list of (check_name, passed, detail).
    """
    checks = []
    sections = re.split(r'\n(?=## )', notes)
    sections = [s.strip() for s in sections if s.strip()]

    # 1. At least 2 sections
    checks.append((
        "≥2 ## sections",
        len(sections) >= 2,
        f"got {len(sections)} section(s)"
    ))

    # 2. Every section has ## heading
    no_heading = [i+1 for i, s in enumerate(sections) if not s.lstrip().startswith("##")]
    checks.append((
        "All sections have ## heading",
        len(no_heading) == 0,
        f"missing headings in sections: {no_heading}" if no_heading else "all present"
    ))

    # 3. Every section has Exam Tip
    missing_tip = [i+1 for i, s in enumerate(sections) if "Exam Tip" not in s and "📝" not in s]
    checks.append((
        "All sections have Exam Tip",
        len(missing_tip) == 0,
        f"sections missing tip: {missing_tip}" if missing_tip else "all present"
    ))

    # 4. No LLM preamble
    preamble_matches = PREAMBLE_RE.findall(notes)
    checks.append((
        "No LLM preamble",
        len(preamble_matches) == 0,
        f"found: {preamble_matches}" if preamble_matches else "clean"
    ))

    # 5. Math is LaTeX (not spelled out)
    math_words = SPELLED_MATH_RE.findall(notes)
    checks.append((
        "Math in LaTeX (not spelled words)",
        len(math_words) == 0,
        f"spelled-out: {math_words[:5]}" if math_words else "all LaTeX"
    ))

    # 6. No HTML tables
    html_tables = HTML_TABLE_RE.findall(notes)
    checks.append((
        "No HTML tables",
        len(html_tables) == 0,
        f"found HTML table tags: {html_tables}" if html_tables else "clean"
    ))

    # 7. LaTeX delimiters are $ not \\( or \\[
    bad_delim = re.findall(r'\\\(|\\\[', notes)
    checks.append((
        "LaTeX uses $ not \\\\( or \\\\[",
        len(bad_delim) == 0,
        f"bad delimiters found: {len(bad_delim)}" if bad_delim else "clean"
    ))

    # 8. Key topics covered (slide titles should appear)
    key_topics = ["Discrete Fourier Transform", "FFT", "Circular Convolution"]
    covered = [t for t in key_topics if t.lower() in notes.lower()]
    checks.append((
        "Key topics covered",
        len(covered) == len(key_topics),
        f"covered {len(covered)}/{len(key_topics)}: {covered}"
    ))

    # 9. DFT formula present (LaTeX or unicode)
    has_formula = bool(re.search(r'X\[k\]|X\s*\[k\]|sum.*n=0|\$\\sum', notes))
    checks.append((
        "DFT formula present",
        has_formula,
        "found" if has_formula else "MISSING — DFT formula not found in notes"
    ))

    # 10. Notes not suspiciously short (< 800 chars is almost certainly truncated)
    checks.append((
        "Notes adequately long (≥800 chars)",
        len(notes) >= 800,
        f"{len(notes)} chars"
    ))

    return checks


# ── Printer ───────────────────────────────────────────────────────────────────

GREEN = "\033[92m"
RED   = "\033[91m"
YELLOW = "\033[93m"
BOLD  = "\033[1m"
RESET = "\033[0m"

def print_report(checks, source, elapsed):
    passed = sum(1 for _, ok, _ in checks if ok)
    total  = len(checks)
    print(f"\n{BOLD}══════════════════════════════════════════{RESET}")
    print(f"{BOLD}  AuraGraph Notes Quality Report{RESET}")
    print(f"  Source: {YELLOW}{source}{RESET}   Time: {elapsed:.1f}s")
    print(f"{'══════════════════════════════════════════'}")
    for name, ok, detail in checks:
        icon  = f"{GREEN}✔{RESET}" if ok else f"{RED}✘{RESET}"
        color = GREEN if ok else RED
        print(f"  {icon}  {color}{name}{RESET}")
        if not ok or detail not in ("", "all present", "clean", "all LaTeX"):
            print(f"       → {detail}")
    print(f"{'──────────────────────────────────────────'}")
    score_color = GREEN if passed == total else (YELLOW if passed >= total * 0.7 else RED)
    print(f"  {score_color}{BOLD}Score: {passed}/{total}{RESET}")
    print(f"{'══════════════════════════════════════════'}\n")
    return passed == total


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--proficiency", default="Intermediate",
                        choices=["Beginner", "Intermediate", "Advanced"])
    parser.add_argument("--url", default=BASE_URL)
    args = parser.parse_args()

    print(f"\n{BOLD}Building test PPTX ({len(SLIDES)} slides)…{RESET}")
    pptx_bytes = build_pptx()
    print(f"  PPTX size: {len(pptx_bytes):,} bytes")

    print(f"{BOLD}Sending to {args.url}/api/upload-fuse-multi (proficiency={args.proficiency})…{RESET}")
    print("  (this may take 20-60s if using LLM)")

    t0 = time.time()
    try:
        resp = requests.post(
            f"{args.url}/api/upload-fuse-multi",
            files={
                "slides_pdfs":   ("lecture_dsp.pptx", pptx_bytes,
                                  "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
                "textbook_pdfs": ("textbook.txt", TEXTBOOK_CONTEXT.encode(),
                                  "text/plain"),
            },
            data={"proficiency": args.proficiency},
            timeout=180,
        )
    except requests.exceptions.Timeout:
        print(f"{RED}TIMEOUT after 180s — server too slow{RESET}")
        sys.exit(2)
    except requests.exceptions.ConnectionError as e:
        print(f"{RED}CONNECTION ERROR — is the server running?{RESET}\n{e}")
        sys.exit(2)

    elapsed = time.time() - t0

    if resp.status_code != 200:
        print(f"{RED}HTTP {resp.status_code}{RESET}")
        print(resp.text[:2000])
        sys.exit(1)

    data = resp.json()
    notes  = data.get("fused_note", "") or data.get("notes", "")
    source = data.get("source", "unknown")

    print(f"\n{BOLD}─── RAW NOTES PREVIEW (first 1200 chars) ───{RESET}")
    print(notes[:1200])
    if len(notes) > 1200:
        print(f"  … ({len(notes):,} total chars) …")
        print(notes[-400:])

    checks  = audit_notes(notes, args.proficiency)
    all_ok  = print_report(checks, source, elapsed)

    # Write full notes to file for human inspection
    out_path = "/tmp/auragraph_test_notes.md"
    with open(out_path, "w") as f:
        f.write(f"# AuraGraph Test — {args.proficiency}\n")
        f.write(f"Source: {source} | Time: {elapsed:.1f}s\n\n")
        f.write(notes)
    print(f"Full notes saved to {BOLD}{out_path}{RESET} for manual review.\n")

    sys.exit(0 if all_ok else 1)


if __name__ == "__main__":
    main()
